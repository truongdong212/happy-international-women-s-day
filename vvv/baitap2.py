from pptx import Presentation
from pptx.util import Inches

# Tạo một bản trình chiếu mới
prs = Presentation()

# Slide 1: Tiêu đề
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sự thắng thế của Thế tử Tông\ntrong 'Kiêu binh nổi loạn'"
subtitle.text = "Người trình bày: (Điền tên bạn)"

# Slide 2: Mở đầu
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Mở đầu"
content.text = "- Giới thiệu văn bản 'Kiêu binh nổi loạn'.\n- Thế tử Tông được lập lên với sự thao túng của kiêu binh."

# Slide 3: Bối cảnh sự kiện
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Bối cảnh sự kiện"
content.text = "- Anh em Quận Huy bị giết, kiêu binh kiểm soát tình hình.\n- Kiêu binh vui mừng, lập Thế tử Tông lên ngôi."

# Slide 4: Kiêu binh rước Thế tử Tông
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Kiêu binh rước Thế tử Tông"
content.text = "- Hô vang, nâng Thế tử lên như một biểu tượng.\n- “Xin ngồi cao thêm nữa…”, “giỡn quả cầu, rước pho tượng Phật.”"

# Slide 5: Ảo tưởng thái bình
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Ảo tưởng thái bình"
content.text = "- Thời tiết thay đổi: “Mấy hôm trước trời u ám, hôm ấy trời trong sáng.”\n- Kiêu binh coi đây là điềm lành, nhưng thực tế chỉ là ảo tưởng."

# Slide 6: Kiêu binh lộng hành
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Kiêu binh lộng hành"
content.text = "- Kiêu binh phóng hỏa, phá nhà Quận Huy.\n- Thế tử Tông không kiểm soát được tình hình."

# Slide 7: Kết luận
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Kết luận"
content.text = "- Thế tử Tông chỉ là bù nhìn, quyền lực thực sự thuộc về kiêu binh.\n- Phản ánh tình trạng rối ren, hỗn loạn của chính quyền thời kỳ này."

# Lưu file
pptx_path = "/mnt/data/Su_Thang_The_The_Tu_Tong.pptx"
prs.save(pptx_path)

pptx_path